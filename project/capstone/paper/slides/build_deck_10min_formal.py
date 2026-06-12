"""Build TALK_10MIN_formal.pptx — the ~10-minute FORMAL cut (academic register).

Same 11-slide structure, figures, numbers, and claims as build_deck_10min.py; titles
and speaker notes are in the register of PAPER_DRAFT.md (no direct address /
contractions / colloquial landings). Intended as a complete base to pare back.
`TALK_10MIN_formal.md` is the readable script. The voice and 8-minute decks are
untouched. Uses the formal geometry figure (slide9_anisotropy_formal.png).

Layout helpers mirror the other builders (self-contained). Run from project/capstone/
(after render_slide_assets.py, which now also emits the formal geometry figure):
    python paper/slides/build_deck_10min_formal.py            # academic theme (default)
    python paper/slides/build_deck_10min_formal.py warm       # -> TALK_10MIN_formal_warm.pptx
    python paper/slides/build_deck_10min_formal.py slate      # -> TALK_10MIN_formal_slate.pptx

Themes are all light (the figures are white-background PNGs; a dark theme would leave
them floating as white cards). All variants import 1:1 into Google Slides, notes included.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FIGS = Path("iteration4_scaling/figures")
ASSETS = Path("paper/slides/assets")

THEMES = {
    # bg=None -> default white, no fill set (cleanest import)
    "academic": dict(bg=None, ink=RGBColor(0x1A, 0x1A, 0x2E), accent=RGBColor(0xC4, 0x4E, 0x52),
                     header=RGBColor(0x4C, 0x72, 0xB0), grey=RGBColor(0x66, 0x66, 0x66),
                     title_font=None, body_font=None, suffix=""),
    # warm journal: off-white paper, serif titles, burgundy/deep-green
    "warm": dict(bg=RGBColor(0xFB, 0xF8, 0xF2), ink=RGBColor(0x2B, 0x2B, 0x28), accent=RGBColor(0x8C, 0x2D, 0x2D),
                 header=RGBColor(0x3A, 0x63, 0x51), grey=RGBColor(0x6E, 0x6A, 0x62),
                 title_font="Georgia", body_font="Calibri", suffix="_warm"),
    # cool slate: pale cool-grey, sans throughout, navy/teal
    "slate": dict(bg=RGBColor(0xF4, 0xF6, 0xFA), ink=RGBColor(0x1F, 0x29, 0x33), accent=RGBColor(0x0F, 0x60, 0x9B),
                  header=RGBColor(0x33, 0x4E, 0x68), grey=RGBColor(0x62, 0x6E, 0x7B),
                  title_font="Segoe UI", body_font="Segoe UI", suffix="_slate"),
}
THEME = THEMES[sys.argv[1] if len(sys.argv) > 1 else "academic"]
OUT = Path(f"paper/slides/TALK_10MIN_formal{THEME['suffix']}.pptx")

DARK = THEME["ink"]
ACCENT = THEME["accent"]
BLUE = THEME["header"]
GREY = THEME["grey"]


def paint(s):
    if THEME["bg"] is not None:
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = THEME["bg"]
    return s

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(title, time_budget):
    s = paint(prs.slides.add_slide(BLANK))
    tb = s.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(11.4), Inches(0.95))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = DARK
    if THEME["title_font"]: p.font.name = THEME["title_font"]
    badge = s.shapes.add_textbox(Inches(12.0), Inches(0.30), Inches(1.1), Inches(0.4))
    bp = badge.text_frame.paragraphs[0]
    bp.text = time_budget
    bp.font.size = Pt(11); bp.font.color.rgb = GREY; bp.alignment = PP_ALIGN.RIGHT
    return s


def bullets(s, items, left=0.6, top=1.35, width=12.1, height=5.6, size=18):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text, kwargs = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(kwargs.get("size", size))
        p.font.bold = kwargs.get("bold", False)
        p.font.italic = kwargs.get("italic", False)
        p.font.color.rgb = kwargs.get("color", DARK)
        if THEME["body_font"]: p.font.name = THEME["body_font"]
        p.level = kwargs.get("level", 0)
        p.space_after = Pt(kwargs.get("space", 10))
    return tb


def image(s, path, left, top, width=None, height=None):
    kw = {}
    if width: kw["width"] = Inches(width)
    if height: kw["height"] = Inches(height)
    return s.shapes.add_picture(str(path), Inches(left), Inches(top), **kw)


def table(s, rows, left, top, width, height, col_widths=None, header_fill=BLUE, size=14):
    shape = s.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top),
                               Inches(width), Inches(height))
    t = shape.table
    if col_widths:
        for j, w in enumerate(col_widths):
            t.columns[j].width = Inches(w)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = t.cell(i, j)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(size); p.font.bold = i == 0
            if THEME["body_font"]: p.font.name = THEME["body_font"]
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return shape


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ---- 1. Title ----------------------------------------------------------------
s = paint(prs.slides.add_slide(BLANK))
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.7), Inches(2.2))
p = tb.text_frame.paragraphs[0]
p.text = "Does relational structure improve linear-probe"
p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = DARK
if THEME["title_font"]: p.font.name = THEME["title_font"]
p2 = tb.text_frame.add_paragraph()
p2.text = "detection of related AI-safety concepts?"
p2.font.size = Pt(34); p2.font.bold = True; p2.font.color.rgb = DARK
if THEME["title_font"]: p2.font.name = THEME["title_font"]
tb2 = s.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.4))
p = tb2.text_frame.paragraphs[0]
p.text = "Jason Boudville"
p.font.size = Pt(22); p.font.color.rgb = GREY
p2 = tb2.text_frame.add_paragraph()
p2.text = "Four iterations, two models, independently replicated — a bounded result."
p2.font.size = Pt(18); p2.font.italic = True; p2.font.color.rgb = ACCENT
notes(s, "This talk asks whether relational structure between concepts improves linear-probe "
         "detection of related AI-safety concepts. Across four iterations — every powered run "
         "pre-registered, with independent peer replication on a second model — we reach a "
         "bounded result: relations help, but only as local hard-negative contrast; the graph's "
         "connectivity contributes no measurable benefit. (0:15)")

# ---- 2. Motivation -----------------------------------------------------------
s = slide("Motivation", "0:55")
image(s, ASSETS / "slide2_graph.png", 0.5, 1.25, width=8.2)
bullets(s, [
    ("Safety monitoring needs a detector — a lens — per behaviour of concern.", {"size": 17}),
    ("These concepts are not independent: they form a relational graph.", {"size": 17}),
    ("Hypothesis: relational structure can be exploited —", {"size": 17, "bold": True}),
    ("    train each lens against its relatives; use connectivity to reduce data.", {"size": 16, "level": 1}),
    ("Relational structure is cheap to specify; labelled data is not.", {"size": 16, "italic": True, "color": ACCENT}),
], left=8.9, top=1.7, width=4.1, height=4.7)
notes(s, "Deployable safety monitoring requires a detector — a lens — for each behaviour of "
         "concern. These concepts are not independent: manipulation, deception, and sycophancy "
         "form a graph of overlapping, related categories. The hypothesis under test is that "
         "this relational structure can be exploited — that training each lens against its "
         "relatives, and exploiting connectivity between concepts, improves detection and "
         "reduces the data required. The appeal is straightforward: relational structure is "
         "inexpensive to specify, whereas labelled data is costly. (0:55)")

# ---- 3. HatCat ---------------------------------------------------------------
s = slide("The HatCat programme and the relational hypothesis", "1:15")
bullets(s, [
    ("HatCat — an interpretability programme building a navigable, interpretable map of the concepts a model represents.",
     {"size": 19, "bold": True}),
    ("", {"size": 8}),
    ("Not a single probe — an atlas of per-concept lenses, connected by a relational graph.", {"size": 18}),
    ("Premise: relational knowledge improves —", {"size": 18}),
    ("    · discrimination (training against neighbouring concepts)", {"size": 17, "level": 1}),
    ("    · coverage efficiency (the graph indicates where to probe; connectivity reaches", {"size": 17, "level": 1}),
    ("      concepts lacking direct data)", {"size": 17, "level": 1}),
    ("", {"size": 8}),
    ("This capstone evaluates the sharpest, most falsifiable form of that premise:", {"size": 18}),
    ("does the graph measurably help, and does connectivity scale?",
     {"size": 20, "bold": True, "color": ACCENT}),
], left=0.7, top=1.35, width=12.0, height=5.7)
notes(s, "The hypothesis originates in HatCat, an interpretability programme whose objective is "
         "a navigable, interpretable map of the concepts a model represents — not a single "
         "probe, but an atlas of per-concept lenses connected by a relational graph. HatCat's "
         "central premise is that relational knowledge improves both discrimination, through "
         "training against neighbouring concepts, and coverage efficiency, since the graph "
         "indicates where to probe and connectivity provides access to concepts that lack "
         "direct data. This capstone evaluates the sharpest, most falsifiable form of that "
         "premise: does relational structure measurably improve a linear probe, and does the "
         "benefit scale with graph connectivity? (1:15)")

# ---- 4. Experimental setup ---------------------------------------------------
s = slide("Experimental setup", "0:30")
image(s, FIGS / "fig1_conditions.png", 0.5, 1.3, width=7.9)
bullets(s, [
    ("15 AI-safety concepts · curated 26-edge graph", {"size": 16}),
    ("4 surface-distinct contexts per concept", {"size": 16}),
    ("Out-of-distribution: train on 3 contexts, test on the held-out 1", {"size": 16, "bold": True}),
    ("Layer-18 mean-pooled residual stream · logistic-regression probes", {"size": 16}),
    ("Gemma-2-9B (+ Gemma-4-E4B peer replication) · 30 seeds", {"size": 16}),
], left=8.6, top=1.7, width=4.4, height=4.5)
notes(s, "The setup follows standard interpretability practice: the mean-pooled residual stream "
         "at layer 18 of Gemma-2-9B, logistic-regression probes, 30 seeds, over 15 concepts and "
         "a curated 26-edge graph. The critical design decision is out-of-distribution "
         "evaluation. Passages are generated across four surface-distinct contexts; probes are "
         "trained on three and tested on the held-out fourth. Without this control, surface "
         "vocabulary saturates discrimination to near-perfect accuracy and the experiment "
         "retains no headroom. (0:30)")

# ---- 5. Iterations 1-2 -------------------------------------------------------
s = slide("Iterations 1–2: a powered, pre-registered null", "1:00")
image(s, ASSETS / "slide4_timeline.png", 0.7, 1.6, width=11.9)
notes(s, "Iteration 1 initially indicated a positive effect. Adversarial verification "
         "identified it as a noise artifact, compounded by a mis-specification: relational "
         "knowledge had been operationalised as additional positive examples rather than as "
         "contrastive structure. The design was rebuilt. Iteration 2 was a pair-matched, "
         "powered test over 105 pairs, with the decision rule fixed in advance. It returned a "
         "clean null — a partial Spearman correlation of minus 0.125 between connectivity and "
         "probe benefit. Because the analysis was pre-registered, the null is reported as a "
         "result. It is subsequently explained, rather than merely recorded. (1:00)")

# ---- 6. Iteration 3 ----------------------------------------------------------
s = slide("Iteration 3: relations as hard negatives", "1:10")
table(s, [
    ["ΔF1 vs random negatives", "E4B (peer replication)", "Gemma-2-9B (this work)"],
    ["direct-neighbour negatives", "+0.057", "+0.059"],
    ["graph-aware mixed", "+0.043", "+0.039"],
    ["held-out-edge (connectivity)", "+0.008", "+0.012"],
], 0.7, 1.6, 8.3, 2.3, col_widths=[3.6, 2.5, 2.2])
bullets(s, [
    ("Path-2 controls (Gemma-2-9B):", {"bold": True, "size": 15}),
    ("curated graph beats placebo permutations, p = 0.002", {"size": 15}),
    ("~⅓ of the effect is shared vocabulary (upper bound)", {"size": 15}),
], left=9.3, top=1.7, width=3.6, height=2.2)
bullets(s, [
    ("Mechanism: contrast, not coverage — each probe is trained against its graph siblings as negatives.",
     {"size": 17, "bold": True}),
    ("Replicated across two models and two codebases.", {"size": 16}),
    ("Bounds: ~⅓ lexical; the benefit is confined to DIRECT adjacency.", {"size": 16, "color": ACCENT}),
], left=0.7, top=4.3, width=12.0, height=2.4)
notes(s, "The correct operationalisation uses related concepts as hard negatives rather than "
         "additional positives — contrast, not coverage. Training each concept's probe against "
         "its graph siblings improves discrimination by approximately 0.06 F1. A placebo-graph "
         "permutation confirms the effect is specific to the curated relationships, at p equals "
         "0.002, rather than an artifact of any graph. Two bounds apply. Approximately one third "
         "of the effect is attributable to shared vocabulary — an upper bound, since the "
         "vocabulary regression over-controls. And the benefit is confined to direct adjacency. "
         "The pattern replicates across both models and codebases. (1:10)")

# ---- 7. Iteration 4 ----------------------------------------------------------
s = slide("Iteration 4: does the connectivity benefit scale?", "1:15")
fig5 = FIGS / "fig5_sweep_with_e4b_overlay.png"
image(s, fig5 if fig5.exists() else FIGS / "fig3_volume_curve.png", 0.8, 1.35, width=8.6)
bullets(s, [
    ("held-out-edge: graph-aware negatives EXCLUDING the evaluated sibling", {"size": 15}),
    ("(the clean alternate-path / Menger test)", {"size": 14, "italic": True, "color": GREY}),
    ("Pre-registered 10× within-model scaling sweep:", {"size": 16, "bold": True}),
    ("top volume +0.005, 95% CI [−0.005, +0.014] — spans zero", {"size": 15, "bold": True}),
    ("minimum detectable effect 0.014; direct effect resolved at 3×", {"size": 15}),
    ("Result: pre-registered NULL — alternate-path prediction refuted at power", {"size": 16, "bold": True, "color": ACCENT}),
    ("E4B = peer-replication corroboration; headline trend is within-model (Gemma-2-9B)",
     {"size": 12, "italic": True, "color": GREY}),
], left=9.5, top=1.7, width=3.5, height=5.2)
notes(s, "The clean test of connectivity, as distinct from adjacency, is the held-out-edge "
         "condition: the probe is trained on all graph-aware negatives except the sibling "
         "against which it is evaluated. If connectivity carries transferable signal, the "
         "remaining relationships should partially reconstruct the withheld boundary. They do "
         "not — the effect is approximately 0.01. To address the possibility that this reflects "
         "insufficient data, we pre-registered a within-model scaling sweep to ten times the "
         "original volume. The held-out-edge delta does not increase; at the top volume it is "
         "0.005, with a 95% confidence interval spanning zero. The minimum detectable effect at "
         "this power was 0.014, and the same design resolves the direct-neighbour effect at "
         "three times that magnitude. The alternate-path, or Menger, prediction is therefore "
         "refuted at adequate power. (1:15)")

# ---- 8. Mechanism ------------------------------------------------------------
s = slide("Mechanism: independent-constraint boundary coverage", "0:35")
image(s, ASSETS / "slide7_cartoon.png", 1.6, 1.4, width=10.0)
notes(s, "A single model accounts for the full pattern. A concept is represented as a centroid "
         "with a surrounding distribution; the probe's boundary lies at the distribution's "
         "tail. Positive examples sharpen the centroid but do not relocate the boundary; each "
         "negative constrains the boundary facet it faces. This predicts that additional "
         "positives have no effect once the centroid is well estimated, that sibling negatives "
         "improve the tested boundary, and that the benefit decreases with distance from it. The "
         "held-out-edge effect is null because, in this graph, the alternate paths route through "
         "the ManipulativeCommunication hub: the constraints are correlated rather than "
         "independent. Connectivity matters only insofar as it counts independent "
         "constraints. (0:35)")

# ---- 9. Geometry control -----------------------------------------------------
s = slide("Geometry control: the anisotropy correction", "0:55")
image(s, ASSETS / "slide9_anisotropy_formal.png", 1.4, 1.55, width=10.5)
notes(s, "An unsupervised geometric check both corroborates the mechanism and illustrates the "
         "analysis discipline. Raw cosine similarity between concept centroids suggested "
         "near-total overlap — 0.971 for adjacent versus 0.961 for non-adjacent pairs — which "
         "would imply the graph is not represented in the model. This is an anisotropy artifact: "
         "transformer activations occupy a narrow cone, inflating all similarities. After "
         "mean-centring, the structure is recovered: adjacent pairs separate by 0.166, "
         "far-domain controls fall well outside the cluster at minus 0.147, while graded graph "
         "distance remains uncorrelated with representation. The geometry independently "
         "reproduces the probe result — local structure is real, multi-hop structure is absent — "
         "and the probe results themselves are unaffected, since a linear classifier removes the "
         "common direction. This is the second instance in which a raw metric was misleading and "
         "a control proved decisive. (0:55)")

# ---- 10. Deployment ----------------------------------------------------------
s = slide("Deployment: rejection scaffolds and tiered negatives", "0:55")
table(s, [
    ["FPR @ 95% TPR", "neighbours", "off-graph confusers", "far-OOD"],
    ["random", "0.19", "0.27", "0.18"],
    ["graph siblings (cheap)", "0.05", "0.26", "0.19"],
    ["model-mined (expensive)", "0.16", "0.06", "0.24"],
    ["mined + cheap “other” tier", "0.15", "0.07", "0.00"],
], 0.7, 1.5, 9.6, 2.9, col_widths=[3.6, 2.0, 2.4, 1.6])
bullets(s, [
    ("FPR @ 95% TPR = false-positive rate at a 95%-true-positive operating point",
     {"size": 13, "italic": True, "color": GREY}),
    ("A flat 15-way classifier has no rejection option; a coarse “other” class rejects out-of-family inputs at full recall.",
     {"size": 16}),
    ("Tiered recipe: inexpensive negatives throughout → audit the off-graph gap → mine only that gap (20-point FP reduction).",
     {"size": 16, "bold": True}),
], left=0.7, top=4.7, width=12.0, height=2.2)
notes(s, "The findings yield a practical recipe. A flat 15-way classifier has no rejection "
         "option and assigns every out-of-domain input to some concept; a coarse 'other' class "
         "rejects out-of-family inputs at full recall with no false rejection of genuine "
         "concepts. With negative count held fixed, false-positive rate at a 95%-true-positive "
         "operating point shows that each scheme seals only the faces it trains against. Graph "
         "siblings seal the neighbour face but not off-graph confusers, against which they "
         "perform no better than random. Model-mined confusers seal that face — a 20-point "
         "reduction — but not the others. And a cheap 'other' tier seals the "
         "far-out-of-distribution face entirely. The recommended pipeline is therefore tiered: "
         "inexpensive negatives throughout, an audit to locate the off-graph gap, and expensive "
         "mining applied only to that gap. Calibration drift is a material failure mode: "
         "operating points set on training data degrade substantially under distribution "
         "shift. (0:55)")

# ---- 11. Conclusions ---------------------------------------------------------
s = slide("Conclusions and future work", "0:50")
bullets(s, [
    ("Relational structure helps only as local hard-negative contrast (~0.06 F1, ~⅓ lexical, two models).",
     {"size": 18, "bold": True}),
    ("The graph's connectivity contributes no measurable benefit at any data volume tested.",
     {"size": 18, "bold": True}),
    ("Pre-registration and adversarial verification were each decisive — twice a raw signal was misleading.",
     {"size": 17}),
    ("The relational structure is a map of where to extend coverage, not a classifier —", {"size": 17, "color": ACCENT}),
    ("    the cartographic objective of HatCat.", {"size": 16, "level": 1, "color": ACCENT}),
    ("", {"size": 8}),
    ("Future work: contrast-prompted generation · near-but-out-of-set rejection test · counting independent constraints.",
     {"size": 15, "color": GREY}),
], top=1.5)
notes(s, "In summary: relational structure improves safety-concept probes only as local "
         "hard-negative contrast — approximately 0.06 F1, roughly one third lexical, replicated "
         "across two models — while the graph's connectivity contributes no measurable benefit "
         "at any data volume tested. Methodologically, pre-registration and adversarial "
         "verification were each decisive: on two occasions a raw signal was misleading and was "
         "identified by a control. Finally, the relational structure is better understood not as "
         "a classifier but as a map of where to extend coverage — which is the cartographic "
         "objective of the HatCat programme. Future work includes contrast-prompted generation, "
         "a near-but-out-of-set rejection test, and a connectivity test that counts independent "
         "constraints rather than raw edges. (0:50)")

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"built {OUT} with {len(prs.slides._sldIdLst)} slides")
