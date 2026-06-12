"""Build TALK_8MIN_voice.pptx — the VOICE-PASS deck (parallel to build_deck.py).

Same slides, same figures, same numbers as build_deck.py; the speaker notes are
reworked toward Jason's spoken voice (dynamic range — plain delivery, spikes at four
beats), the title is plain, and a few on-slide framing lines are tightened. The layout
helpers mirror build_deck.py deliberately (kept self-contained so the two decks stay
independent for side-by-side comparison; only content/notes differ). `TALK_8MIN_voice.md`
is the readable script.

Run from project/capstone/ (after render_slide_assets.py):
    python paper/slides/build_deck_voice.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FIGS = Path("iteration4_scaling/figures")
ASSETS = Path("paper/slides/assets")
OUT = Path("paper/slides/TALK_8MIN_voice.pptx")

DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xC4, 0x4E, 0x52)
BLUE = RGBColor(0x4C, 0x72, 0xB0)
GREY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(title, time_budget):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(11.4), Inches(0.95))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = DARK
    badge = s.shapes.add_textbox(Inches(12.0), Inches(0.30), Inches(1.1), Inches(0.4))
    bp = badge.text_frame.paragraphs[0]
    bp.text = time_budget
    bp.font.size = Pt(11); bp.font.color.rgb = GREY; bp.alignment = PP_ALIGN.RIGHT
    return s


def bullets(s, items, left=0.6, top=1.35, width=12.1, height=5.6, size=18):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text, kwargs = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(kwargs.get("size", size))
        p.font.bold = kwargs.get("bold", False)
        p.font.italic = kwargs.get("italic", False)
        p.font.color.rgb = kwargs.get("color", DARK)
        p.level = kwargs.get("level", 0)
        p.space_after = Pt(kwargs.get("space", 10))
    return tb


def image(s, path, left, top, width=None, height=None):
    kw = {}
    if width: kw["width"] = Inches(width)
    if height: kw["height"] = Inches(height)
    return s.shapes.add_picture(str(path), Inches(left), Inches(top), **kw)


def table(s, rows, left, top, width, height, col_widths=None, header_fill=BLUE, size=14):
    shape = s.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top),
                               Inches(width), Inches(height))
    t = shape.table
    if col_widths:
        for j, w in enumerate(col_widths):
            t.columns[j].width = Inches(w)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = t.cell(i, j)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(size); p.font.bold = i == 0
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return shape


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ---- 1. Title ----------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.8))
p = tb.text_frame.paragraphs[0]
p.text = "Does relational structure improve"
p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = DARK
p2 = tb.text_frame.add_paragraph()
p2.text = "safety-concept probes?"
p2.font.size = Pt(40); p2.font.bold = True; p2.font.color.rgb = DARK
tb2 = s.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(11.7), Inches(1.4))
p = tb2.text_frame.paragraphs[0]
p.text = "Jason Boudville"
p.font.size = Pt(22); p.font.color.rgb = GREY
p2 = tb2.text_frame.add_paragraph()
p2.text = "Four iterations, two models, independently replicated — and a bounded answer."
p2.font.size = Pt(18); p2.font.italic = True; p2.font.color.rgb = ACCENT
notes(s, "Here's the short version, so you know where we're headed. We took an idea that sounds "
         "obviously true, and we tested it hard enough to trust the answer — including the part "
         "of the answer that turned out to be no. (0:15)")

# ---- 2. The question ---------------------------------------------------------
s = slide("The question", "1:00")
image(s, ASSETS / "slide2_graph.png", 0.5, 1.25, width=8.2)
bullets(s, [
    ("The hypothesis:", {"bold": True, "size": 17}),
    ("Probes for related concepts should help each other —", {"size": 16}),
    ("and the benefit should scale with connectivity (alternate paths).", {"size": 16}),
    ("Like inferring a border from the neighbours' borders.", {"italic": True, "size": 15, "color": GREY}),
    ("Graphs are cheap. Data is not.", {"bold": True, "size": 16, "color": ACCENT}),
], left=8.9, top=1.8, width=4.1, height=4.5)
notes(s, "Safety monitoring wants a detector for each behaviour we care about — a lens for "
         "manipulation, one for deception, one for sycophancy. But these concepts aren't "
         "strangers to each other. They overlap, they blur at the edges, they form a graph. So "
         "here's the idea we inherited — and someone later checked it on a second model. You "
         "lean on that graph. You train each lens against its relatives. And where two concepts "
         "are joined by lots of independent paths, you pin the boundary between them without "
         "ever showing the probe that boundary directly. It's the appealing kind of idea. "
         "Structure is free. Data isn't. (1:00)")

# ---- 3. Setup ----------------------------------------------------------------
s = slide("Setup in 30 seconds", "0:35")
image(s, FIGS / "fig1_conditions.png", 0.5, 1.3, width=7.9)
bullets(s, [
    ("15 AI-safety concepts · 26-edge curated graph", {"size": 16}),
    ("Passages in 4 surface-distinct contexts", {"size": 16}),
    ("Train on 3 contexts, test on the held-out one (OOD)", {"size": 16, "bold": True}),
    ("Linear probes, layer-18 residual stream", {"size": 16}),
    ("Gemma-2-9B (+ independent E4B replication) · 30 seeds", {"size": 16}),
], left=8.6, top=1.7, width=4.4, height=4.5)
notes(s, "Thirty seconds of setup, all of it the ARENA toolkit. Mean-pooled residual stream, "
         "logistic regression on top. The one choice that matters: we always test out of "
         "distribution — train the probe in three settings, test it in a fourth it's never "
         "seen. Skip that and surface vocabulary does all the work, every probe scores perfect, "
         "and you've measured nothing. (0:35)")

# ---- 4. Act 1: the honest null -----------------------------------------------
s = slide("Act 1 — the honest null", "1:00")
image(s, ASSETS / "slide4_timeline.png", 0.7, 1.6, width=11.9)
notes(s, "Iteration one looked positive. It was wrong. We pulled it apart and found noise, not "
         "signal — and worse, we'd asked the wrong question, feeding the probe more examples of "
         "each concept instead of teaching it the boundary. So we "
         "rebuilt it properly. Iteration two: a hundred and five matched pairs, powered, the "
         "decision rule locked before we looked. The answer came back a clean null — "
         "connectivity does not predict how much a probe improves. Because we'd pre-registered, "
         "there was nowhere to hide it and no reason to. Hold onto that null. Two slides from "
         "now it stops being a disappointment and becomes the point. (1:00)")

# ---- 5. Act 2: the reframe ---------------------------------------------------
s = slide("Act 2 — the reframe that worked: relations as hard negatives", "1:15")
table(s, [
    ["ΔF1 vs random negatives", "E4B (peer replication)", "Gemma-2-9B (this work)"],
    ["direct-neighbour negatives", "+0.057", "+0.059"],
    ["graph-aware mixed", "+0.043", "+0.039"],
    ["held-out-edge (connectivity)", "+0.008", "+0.012"],
], 0.7, 1.6, 8.3, 2.3, col_widths=[3.6, 2.5, 2.2])
bullets(s, [
    ("Path-2 controls (Gemma-2-9B):", {"bold": True, "size": 15}),
    ("declared graph beats shuffled placebo graphs, p = 0.002", {"size": 15}),
    ("~⅓ of the effect is shared vocabulary", {"size": 15}),
], left=9.3, top=1.7, width=3.6, height=2.2)
bullets(s, [
    ("The mechanism is contrast, not coverage: train each concept's probe with its actual siblings as the negative class.",
     {"size": 17, "bold": True}),
    ("Robust: two models, two codebases, fully independent runs — same numbers.", {"size": 16}),
    ("Honest bounds: ~⅓ lexical, and the effect lives entirely in DIRECT adjacency →", {"size": 16, "color": ACCENT}),
], left=0.7, top=4.3, width=12.0, height=2.4)
notes(s, "Here's what we'd missed. The fix isn't more examples of a concept — it's the right "
         "negatives. Train manipulation's probe against its actual siblings, the concepts it "
         "keeps getting confused with, and discrimination jumps about six points of F1. And the "
         "graph earns its keep: shuffle the edges into a fake graph and the effect vanishes — "
         "p of nought-point-zero-zero-two — so these particular relationships are real, not just "
         "any old structure. Two honest caveats, because that's the job. About a third of the "
         "lift is shared vocabulary, not deep structure. And — this is the column on the right — "
         "all of it lives in direct neighbours. Same numbers, two models, two codebases, run "
         "independently. (1:15)")

# ---- 6. Act 3: connectivity isolated -----------------------------------------
s = slide("Act 3 — connectivity isolated: the pre-registered NULL", "1:15")
fig5 = FIGS / "fig5_sweep_with_e4b_overlay.png"
image(s, fig5 if fig5.exists() else FIGS / "fig3_volume_curve.png", 0.8, 1.35, width=8.6)
bullets(s, [
    ("held-out-edge = train with all graph-aware negatives EXCEPT the tested sibling", {"size": 15}),
    ("If alternate paths carry signal, they should reconstruct the withheld boundary", {"size": 15}),
    ("They don't — at any volume:", {"bold": True, "size": 17, "color": ACCENT}),
    ("10× sweep: top volume +0.005, 95% CI [−0.005, +0.014] — spans zero", {"size": 15, "bold": True}),
    ("Power shown: could detect 0.014; the direct effect is 3× that", {"size": 15}),
    ("VERDICT: pre-registered NULL", {"size": 17, "bold": True, "color": ACCENT}),
    ("E4B points = peer-replication corroboration; headline trend is within-model (Gemma-2-9B)",
     {"size": 12, "italic": True, "color": GREY}),
], left=9.5, top=1.7, width=3.5, height=5.2)
notes(s, "Now the clean test of connectivity, separated from mere adjacency. Train with all the "
         "graph's negatives except the very sibling you're being tested on. If the graph's "
         "deeper structure carries anything, those other relationships should reconstruct the "
         "boundary you held back. They don't — about a hundredth of a point. And the obvious "
         "comeback is 'you just didn't have enough data.' So we settled it. We bought ten times "
         "the data, pre-registered, one model start to finish. The line didn't move. At the top "
         "volume it's five-thousandths of a point, with a confidence interval sitting squarely "
         "across zero — say just that. And before anyone asks whether we had the power to see an "
         "effect: the design could catch one of fourteen-thousandths, and it picks up the direct "
         "effect — three times that — without breaking a sweat. The instrument works. The thing "
         "we were testing just isn't there. (1:15)")

# ---- 7. Why ------------------------------------------------------------------
s = slide("Why — the boundary-coverage picture", "0:35")
image(s, ASSETS / "slide7_cartoon.png", 1.6, 1.4, width=10.0)
notes(s, "One picture explains all of it. A probe's boundary is set by whichever negative sits "
         "closest to it. More examples of the concept — nothing; the centre's already sharp. A "
         "sibling negative — it leans right on the face you're testing, so it helps. And the "
         "held-out-edge stays at zero because in our graph the alternate routes run back through "
         "the same hub — manipulation — so those paths aren't independent. They're the same "
         "information wearing different hats. Menger counted paths. What a boundary actually "
         "needs is independent constraints.\n\n"
         "If time: the direct effect itself shrinks as data grows, +0.061 down to +0.040 — so "
         "smart negatives matter most exactly when data is scarce, which is where a long-tail "
         "safety concept lives. (0:35)")

# ---- 8. So what --------------------------------------------------------------
s = slide("So what — you seal exactly the faces you train against", "1:00")
table(s, [
    ["Negative set (FPR @ 95% TPR)", "neighbours", "off-graph confusers", "far-OOD"],
    ["random", "0.19", "0.27", "0.18"],
    ["graph siblings (cheap)", "0.05", "0.26", "0.19"],
    ["model-mined (expensive)", "0.16", "0.06", "0.24"],
    ["mined + cheap “other” tier", "0.15", "0.07", "0.00"],
], 0.7, 1.5, 9.6, 2.9, col_widths=[3.6, 2.0, 2.4, 1.6])
bullets(s, [
    ("FPR@95%TPR = false-positive rate at an operating point keeping 95% of true positives",
     {"size": 13, "italic": True, "color": GREY}),
    ("A flat 15-probe argmax has no “none of these” — it labels a dog passage as manipulation. A cheap OTHER class fixes that completely.",
     {"size": 16}),
    ("Recipe: cheap tiers everywhere → audit for the gap → spend mining only on the gap (a 20-point FP cut).",
     {"size": 16, "bold": True}),
], left=0.7, top=4.7, width=12.0, height=2.2)
notes(s, "This pays off in something you can actually build. A flat fifteen-way probe has no "
         "'none of these' — show it a passage about a dog and it'll still confidently call it "
         "manipulation. Add one cheap reject class and that problem's gone. After that, picking "
         "negatives is just budget. Graph siblings seal the neighbour face for nothing. A cheap "
         "'other' tier seals the far-away stuff completely. But the graph gives you almost no "
         "protection against the model's real confusers — the ones it actually trips on — and "
         "only the expensive option, mining the model's own confusions, seals those: a "
         "twenty-point cut in false positives. So: cheap everywhere, find the gap, spend the "
         "expensive money only on the gap. (1:00)")

# ---- 9. Close ----------------------------------------------------------------
s = slide("What this taught me", "0:45")
bullets(s, [
    ("Pre-register before powered runs — saved us twice", {"size": 22, "bold": True}),
    ("Adversarially verify your positives — killed a false one", {"size": 22, "bold": True}),
    ("A bounded answer + a mechanism beats an unbounded yes", {"size": 22, "bold": True}),
    ("", {"size": 10}),
    ("Next:", {"size": 17, "color": GREY}),
    ("contrast-prompted generation test · near-out-of-set rejection test (controls generated + audited) · count independent constraints, not edges",
     {"size": 15, "color": GREY}),
], top=1.8)
notes(s, "What I'll take from this. Iteration one's false positive would have made a great talk "
         "— and it would have been wrong. The discipline is the only reason I can stand here and "
         "trust these numbers: pre-register, try to kill your own positives, bound what you "
         "claim. And the relationships? They turned out not to be a way of telling concepts "
         "apart at all. They're a map — of where to look next. Thanks. Happy to take "
         "questions. (0:45)")

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"built {OUT} with {len(prs.slides._sldIdLst)} slides")
