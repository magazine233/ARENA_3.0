"""Build TALK_8MIN.pptx from the scaffold content (paper/slides/TALK_8MIN.md is the
authoritative script; this file is the deck source — edit content here, rebuild).

Run from project/capstone/ (after render_slide_assets.py):
    python paper/slides/build_deck.py

Slide 6 auto-upgrades: if iteration4_scaling/figures/fig5_sweep_with_e4b_overlay.png
exists (exported by ITERATION4_results_explorer.ipynb after STEP 3), it is used;
otherwise the existing 7/10/14 volume curve (fig3) appears with a PLACEHOLDER banner.
Speaker notes (the SAY text) are embedded in each slide's notes pane.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FIGS = Path("iteration4_scaling/figures")
ASSETS = Path("paper/slides/assets")
OUT = Path("paper/slides/TALK_8MIN.pptx")

DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xC4, 0x4E, 0x52)
BLUE = RGBColor(0x4C, 0x72, 0xB0)
GREY = RGBColor(0x66, 0x66, 0x66)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(title, time_budget):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(11.4), Inches(0.95))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = DARK
    badge = s.shapes.add_textbox(Inches(12.0), Inches(0.30), Inches(1.1), Inches(0.4))
    bp = badge.text_frame.paragraphs[0]
    bp.text = time_budget
    bp.font.size = Pt(11)
    bp.font.color.rgb = GREY
    bp.alignment = PP_ALIGN.RIGHT
    return s


def bullets(s, items, left=0.6, top=1.35, width=12.1, height=5.6, size=18):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text, kwargs = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(kwargs.get("size", size))
        p.font.bold = kwargs.get("bold", False)
        p.font.italic = kwargs.get("italic", False)
        p.font.color.rgb = kwargs.get("color", DARK)
        p.level = kwargs.get("level", 0)
        p.space_after = Pt(kwargs.get("space", 10))
    return tb


def image(s, path, left, top, width=None, height=None):
    kw = {}
    if width: kw["width"] = Inches(width)
    if height: kw["height"] = Inches(height)
    return s.shapes.add_picture(str(path), Inches(left), Inches(top), **kw)


def table(s, rows, left, top, width, height, col_widths=None, header_fill=BLUE, size=14):
    shape = s.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top),
                               Inches(width), Inches(height))
    t = shape.table
    if col_widths:
        for j, w in enumerate(col_widths):
            t.columns[j].width = Inches(w)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = t.cell(i, j)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(size)
            p.font.bold = i == 0
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return shape


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ---- 1. Title -----------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.8))
p = tb.text_frame.paragraphs[0]
p.text = "Knowing the Graph:"
p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = DARK
p2 = tb.text_frame.add_paragraph()
p2.text = "does relational structure improve safety-concept probes?"
p2.font.size = Pt(32); p2.font.color.rgb = DARK
tb2 = s.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(11.7), Inches(1.4))
p = tb2.text_frame.paragraphs[0]
p.text = "Jason Boudville"
p.font.size = Pt(22); p.font.color.rgb = GREY
p2 = tb2.text_frame.add_paragraph()
p2.text = "Four iterations, two models, independently replicated — and a bounded answer."
p2.font.size = Pt(18); p2.font.italic = True; p2.font.color.rgb = ACCENT
notes(s, "This is a story about testing an appealing hypothesis carefully enough to trust "
         "the answer — including the parts of the answer that are 'no'. (0:15)")

# ---- 2. The question ------------------------------------------------------------
s = slide("The question", "1:00")
image(s, ASSETS / "slide2_graph.png", 0.5, 1.25, width=8.2)
bullets(s, [
    ("The hypothesis:", {"bold": True, "size": 17}),
    ("Probes for related concepts should help each other —", {"size": 16}),
    ("and the benefit should scale with connectivity (alternate paths).", {"size": 16}),
    ("Like inferring a border from the neighbours' borders.", {"italic": True, "size": 15, "color": GREY}),
    ("Graphs are cheap. Data is not.", {"bold": True, "size": 16, "color": ACCENT}),
], left=8.9, top=1.8, width=4.1, height=4.5)
notes(s, "Safety monitoring wants per-concept detectors — lenses — over families of related "
         "behaviours. These concepts aren't independent: manipulation, deception, sycophancy "
         "form a graph. The hypothesis we inherited — later peer-replicated on a second "
         "model: relationship-aware training improves detection, and if two concepts are "
         "connected by many independent paths, you can triangulate their boundary even without "
         "direct data. Intuitive, testable, and if true, very useful: graphs are cheap, data is "
         "not. (1:00)")

# ---- 3. Setup ------------------------------------------------------------------
s = slide("Setup in 30 seconds", "0:35")
image(s, FIGS / "fig1_conditions.png", 0.5, 1.3, width=7.9)
bullets(s, [
    ("15 AI-safety concepts · 26-edge curated graph", {"size": 16}),
    ("Passages in 4 surface-distinct contexts", {"size": 16}),
    ("Train on 3 contexts, test on the held-out one (OOD)", {"size": 16, "bold": True}),
    ("Linear probes, layer-18 residual stream", {"size": 16}),
    ("Gemma-2-9B (+ independent E4B replication) · 30 seeds", {"size": 16}),
], left=8.6, top=1.7, width=4.4, height=4.5)
notes(s, "Everything is the ARENA toolkit you know: mean-pooled residual stream, logistic "
         "regression. The one design choice to remember: out-of-distribution evaluation — "
         "train on three contexts, test on the fourth — because without it, surface vocabulary "
         "saturates everything and every probe looks perfect. (0:35)")

# ---- 4. Act 1: the honest null ----------------------------------------------------
s = slide("Act 1 — the honest null", "1:00")
image(s, ASSETS / "slide4_timeline.png", 0.7, 1.6, width=11.9)
notes(s, "Iteration 1 looked positive — and died under adversarial verification: a noise "
         "artifact, plus we'd operationalised 'knowing the relations' as extra POSITIVE "
         "passages. Iteration 2 fixed validity and power — 105 matched pairs, decision rule "
         "locked before the run — and the answer was a clean null: connectivity does not "
         "predict probe benefit. We pre-registered, so we report it as a result, not a "
         "failure. Hold that thought, because the null gets EXPLAINED in two slides. (1:00)")

# ---- 5. Act 2: the reframe ------------------------------------------------------
s = slide("Act 2 — the reframe that worked: relations as hard negatives", "1:15")
table(s, [
    ["ΔF1 vs random negatives", "E4B (peer replication)", "Gemma-2-9B (this work)"],
    ["direct-neighbour negatives", "+0.057", "+0.059"],
    ["graph-aware mixed", "+0.043", "+0.039"],
    ["held-out-edge (connectivity)", "+0.008", "+0.012"],
], 0.7, 1.6, 8.3, 2.3, col_widths=[3.6, 2.5, 2.2])
tb = bullets(s, [
    ("Path-2 controls (Gemma-2-9B):", {"bold": True, "size": 15}),
    ("declared graph beats shuffled placebo graphs, p = 0.002", {"size": 15}),
    ("~⅓ of the effect is shared vocabulary", {"size": 15}),
], left=9.3, top=1.7, width=3.6, height=2.2)
bullets(s, [
    ("The mechanism is contrast, not coverage: train each concept's probe with its actual siblings as the negative class.",
     {"size": 17, "bold": True}),
    ("Robust: two models, two codebases, fully independent runs — same numbers.", {"size": 16}),
    ("Honest bounds: ~⅓ lexical, and the effect lives entirely in DIRECT adjacency →", {"size": 16, "color": ACCENT}),
], left=0.7, top=4.3, width=12.0, height=2.4)
notes(s, "The peer replication's diagnosis: the mechanism isn't extra positives, it's hard "
         "negatives — train Manipulation's probe with its actual siblings as the negative "
         "class. That works: +0.06 F1, and the declared graph beats shuffled placebo graphs at "
         "p=0.002 — the curated edges are real confusability structure, not just 'any graph'. "
         "Two honest bounds: about a third of the effect is shared vocabulary, and — right "
         "column — the effect lives entirely in DIRECT adjacency. Same numbers, two models, "
         "two codebases, fully independent runs. (1:15)")

# ---- 6. Act 3: connectivity isolated [ITER-4 SLOT] -------------------------------
s = slide("Act 3 — connectivity isolated: does it scale?", "1:15")
fig5 = FIGS / "fig5_sweep_with_e4b_overlay.png"
if fig5.exists():
    image(s, fig5, 0.8, 1.35, width=8.6)
else:
    image(s, FIGS / "fig3_volume_curve.png", 0.8, 1.7, width=8.2)
    banner = s.shapes.add_textbox(Inches(0.8), Inches(1.30), Inches(8.2), Inches(0.4))
    bp = banner.text_frame.paragraphs[0]
    bp.text = "PLACEHOLDER — final figure (fig5, volumes to 70/context) lands after the iteration-4 sweep"
    bp.font.size = Pt(14); bp.font.bold = True; bp.font.color.rgb = ORANGE
bullets(s, [
    ("held-out-edge = train with all graph-aware negatives EXCEPT the tested sibling", {"size": 15}),
    ("If alternate paths carry signal, they should reconstruct the withheld boundary", {"size": 15}),
    ("They don't: ≈ +0.01", {"bold": True, "size": 18, "color": ACCENT}),
    ("[ITER-4 VERDICT HERE after the sweep]", {"size": 15, "color": ORANGE, "bold": True}),
    ("E4B points = independent corroboration; headline trend is within-model (Gemma-2-9B)",
     {"size": 12, "italic": True, "color": GREY}),
], left=9.5, top=1.7, width=3.5, height=5.0)
notes(s, "The clean test of CONNECTIVITY — as opposed to adjacency — is the held-out-edge "
         "condition: train with all graph-aware negatives EXCEPT the sibling you're tested "
         "against. If alternate paths carry signal, they should partially reconstruct that "
         "withheld boundary. They don't: about +0.01.\n\n"
         "[NULL FILL: And the obvious rebuttal — 'you just need more data' — is what iteration "
         "4 killed: we scaled to 10x under a locked pre-registration, and the line stays flat. "
         "Say aloud only the top-volume mean and CI; per-volume numbers and p-values live on "
         "the slide. Within-model, pre-registered, minimum detectable effect reported: the "
         "alternate-path prediction is refuted with shown — not asserted — power.]\n\n"
         "[CONFIRM FILL: And iteration 4 surprised us: at 10x volume the held-out-edge delta "
         "rises — numbers — connectivity was power-limited, and the next question is WHICH "
         "edges recovered signal.] (1:15)")

# ---- 7. Why ----------------------------------------------------------------------
s = slide("Why — the boundary-coverage picture", "0:35")
image(s, ASSETS / "slide7_cartoon.png", 1.6, 1.4, width=10.0)
notes(s, "One model predicts everything you've seen: a probe's boundary is set by the "
         "negatives nearest each facet. Extra positives — nothing (iterations 1-2). Sibling "
         "negatives — the facet you're tested on (+0.06). Held-out-edge — the surviving "
         "constraints face the wrong way, and in our graph the alternate paths all pass "
         "through the ManipulativeCommunication hub, so they're correlated. Menger counted "
         "paths; what matters is INDEPENDENT constraints.\n\n"
         "[If iteration 4 CONFIRMS instead, close with: '—and iteration 4 showed this model "
         "undersold independence: some alternate paths do carry signal at volume. Finding "
         "which ones is the next sub-analysis.'] (0:35)")

# ---- 8. So what -------------------------------------------------------------------
s = slide("So what — you seal exactly the faces you train against", "1:00")
table(s, [
    ["Negative set (FPR @ 95% TPR)", "neighbours", "off-graph confusers", "far-OOD"],
    ["random", "0.19", "0.27", "0.18"],
    ["graph siblings (cheap)", "0.05", "0.26", "0.19"],
    ["model-mined (expensive)", "0.16", "0.06", "0.24"],
    ["mined + cheap “other” tier", "0.15", "0.07", "0.00"],
], 0.7, 1.5, 9.6, 2.9, col_widths=[3.6, 2.0, 2.4, 1.6])
bullets(s, [
    ("FPR@95%TPR = false-positive rate at an operating point keeping 95% of true positives",
     {"size": 13, "italic": True, "color": GREY}),
    ("A flat 15-probe argmax has no “none of these” — it labels a dog passage as manipulation. A cheap OTHER class fixes that completely.",
     {"size": 16}),
    ("Recipe: cheap tiers everywhere → audit for the gap → spend mining only on the gap (a 20-point FP cut).",
     {"size": 16, "bold": True}),
], left=0.7, top=4.7, width=12.0, height=2.2)
notes(s, "This cashes out practically. A flat 15-probe argmax has no 'none of these' — it "
         "confidently labels a dog passage as manipulation; a cheap OTHER class fixes that "
         "completely. And negative selection is a budget allocation: graph siblings seal the "
         "neighbour face for free; the cheap 'other' tier seals far-OOD; but the graph gives "
         "~zero protection against the model's ACTUAL off-graph confusers — only model-mining "
         "seals those, a 20-point FP cut. So: cheap tiers everywhere, audit for the gap, spend "
         "mining only on the gap. (1:00)")

# ---- 9. Close ----------------------------------------------------------------------
s = slide("What this taught me", "0:45")
bullets(s, [
    ("Pre-register before powered runs — saved us twice", {"size": 22, "bold": True}),
    ("Adversarially verify your positives — killed a false one", {"size": 22, "bold": True}),
    ("A bounded answer + a mechanism beats an unbounded yes", {"size": 22, "bold": True}),
    ("", {"size": 10}),
    ("Next:", {"size": 17, "color": GREY}),
    ("contrast-prompted generation test · near-out-of-set rejection test (controls generated + audited) · count independent constraints, not edges",
     {"size": 15, "color": GREY}),
], top=1.8)
notes(s, "The meta-result for this room: iteration 1's false positive would have been a fun "
         "talk and wrong. The discipline — pre-registration, adversarial verification, honest "
         "bounding — is what let two fully independent runs converge on numbers we trust. The "
         "relationships turned out to matter not as classifiers but as a map: where to look "
         "next. Thanks — questions. (0:45)")

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"built {OUT} ({len(prs.slides.slides if hasattr(prs.slides, 'slides') else prs.slides._sldIdLst)} slides)"
      if False else f"built {OUT} with {len(prs.slides._sldIdLst)} slides")
